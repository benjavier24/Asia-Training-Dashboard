import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r'c:\Users\BenjJavier\OneDrive - bolttech\Documents\Copilot\Created\training_dashboard\pictogram.pptx'
prs = Presentation(path)

colors_found = set()

# Check slide backgrounds and shape fills
for i, slide in enumerate(prs.slides):
    # Check slide background
    bg = slide.background
    fill = bg.fill
    if fill.type is not None:
        try:
            if fill.fore_color and fill.fore_color.rgb:
                colors_found.add(("BG", i+1, str(fill.fore_color.rgb)))
        except:
            pass
    
    # Check shapes for colors
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color and run.font.color.rgb:
                            colors_found.add(("TEXT", i+1, str(run.font.color.rgb)))
                    except:
                        pass
        
        # Check shape fill
        if hasattr(shape, 'fill'):
            try:
                f = shape.fill
                if f.type is not None and f.fore_color and f.fore_color.rgb:
                    colors_found.add(("FILL", i+1, str(f.fore_color.rgb)))
            except:
                pass

# Print unique colors
print("=== COLORS FOUND IN BOLTTECH PICTOGRAM FILE ===\n")

text_colors = sorted(set(c[2] for c in colors_found if c[0] == "TEXT"))
fill_colors = sorted(set(c[2] for c in colors_found if c[0] == "FILL"))
bg_colors = sorted(set(c[2] for c in colors_found if c[0] == "BG"))

print(f"TEXT colors ({len(text_colors)}):")
for c in text_colors:
    print(f"  #{c}")

print(f"\nFILL colors ({len(fill_colors)}):")
for c in fill_colors:
    print(f"  #{c}")

print(f"\nBACKGROUND colors ({len(bg_colors)}):")
for c in bg_colors:
    print(f"  #{c}")

# Also check the slide layout names for theme info
print("\n=== SLIDE LAYOUT THEMES ===")
layouts = set()
for slide in prs.slides:
    if slide.slide_layout:
        layouts.add(slide.slide_layout.name)
for l in sorted(layouts):
    print(f"  {l}")
