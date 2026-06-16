import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

path = r'c:\Users\BenjJavier\OneDrive - bolttech\Documents\Copilot\Created\training_dashboard\pictogram.pptx'
print(f"File exists: {os.path.exists(path)}")
print(f"File size: {os.path.getsize(path)}")

prs = Presentation(path)

print(f"\nTotal slides: {len(prs.slides)}")
print("=" * 80)

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    if slide.slide_layout:
        print(f"  Layout: {slide.slide_layout.name}")
    
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            text = shape.text.strip().replace('\n', ' | ')
            print(f"  TEXT: '{text[:120]}'")
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"  PICTURE: name='{shape.name}'")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            texts = []
            for sh in shape.shapes:
                if sh.has_text_frame and sh.text.strip():
                    texts.append(sh.text.strip())
            if texts:
                print(f"  GROUP: [{', '.join(texts[:5])}]")
            else:
                print(f"  GROUP: (no text, {len(shape.shapes)} sub-shapes)")
        else:
            print(f"  SHAPE: type={shape.shape_type}, name='{shape.name}'")
